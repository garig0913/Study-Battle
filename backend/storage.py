import os
import uuid
import json
import tempfile
from typing import List, Dict, Tuple, Optional
import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
from PIL import Image
import io
import re


PERSIST_BASE = os.path.join(os.path.dirname(__file__), "data")
PERSIST_UPLOADS = os.path.join(PERSIST_BASE, "uploads")
PERSIST_COURSES = os.path.join(PERSIST_BASE, "courses")
UPLOAD_DIR = PERSIST_UPLOADS
MAX_FILE_SIZE = 20 * 1024 * 1024


class TextExtractor:
    
    @staticmethod
    def extract_from_pdf(file_path: str) -> List[Dict]:
        pages = []
        try:
            with pdfplumber.open(file_path) as pdf:
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text() or ""
                    pages.append({
                        "page_number": i + 1,
                        "text": text,
                        "char_start": sum(len(p["text"]) for p in pages),
                    })
        except Exception as e:
            print(f"Error extracting PDF: {e}")
        return pages
    
    @staticmethod
    def extract_from_docx(file_path: str) -> List[Dict]:
        pages = []
        try:
            doc = DocxDocument(file_path)
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            text = "\n".join(full_text)
            pages.append({
                "page_number": 1,
                "text": text,
                "char_start": 0,
            })
        except Exception as e:
            print(f"Error extracting DOCX: {e}")
        return pages
    
    @staticmethod
    def extract_from_pptx(file_path: str) -> List[Dict]:
        pages = []
        try:
            prs = Presentation(file_path)
            char_offset = 0
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                text = "\n".join(slide_text)
                pages.append({
                    "page_number": i + 1,
                    "text": text,
                    "char_start": char_offset,
                })
                char_offset += len(text)
        except Exception as e:
            print(f"Error extracting PPTX: {e}")
        return pages
    
    @staticmethod
    def extract_from_txt(file_path: str) -> List[Dict]:
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
            return [{
                "page_number": 1,
                "text": text,
                "char_start": 0,
            }]
        except Exception as e:
            print(f"Error extracting TXT: {e}")
            return []
    
    @staticmethod
    def extract_from_image(file_path: str) -> List[Dict]:
        return [{
            "page_number": 1,
            "text": "[Image file - OCR not available in this demo]",
            "char_start": 0,
        }]
    
    @classmethod
    def extract(cls, file_path: str, file_name: str) -> List[Dict]:
        ext = os.path.splitext(file_name)[1].lower()
        
        if ext == ".pdf":
            return cls.extract_from_pdf(file_path)
        elif ext == ".docx":
            return cls.extract_from_docx(file_path)
        elif ext == ".pptx":
            return cls.extract_from_pptx(file_path)
        elif ext == ".txt":
            return cls.extract_from_txt(file_path)
        elif ext in [".png", ".jpg", ".jpeg", ".gif", ".bmp"]:
            return cls.extract_from_image(file_path)
        else:
            return cls.extract_from_txt(file_path)


class TextChunker:
    
    def __init__(self, chunk_size: int = 500, overlap: int = 50):
        self.chunk_size = chunk_size
        self.overlap = overlap
    
    def _estimate_tokens(self, text: str) -> int:
        return len(text) // 4
    
    def chunk_text(
        self, 
        text: str, 
        doc_id: str, 
        file_name: str, 
        page_number: int,
        base_char_offset: int = 0
    ) -> List[Dict]:
        chunks = []
        
        text = re.sub(r'\s+', ' ', text).strip()
        
        if not text:
            return chunks
        
        words = text.split(' ')
        current_chunk = []
        current_char_start = base_char_offset
        char_position = base_char_offset
        
        for word in words:
            current_chunk.append(word)
            
            chunk_text = ' '.join(current_chunk)
            if self._estimate_tokens(chunk_text) >= self.chunk_size:
                chunk_id = str(uuid.uuid4())
                chunk_end = char_position + len(word)
                
                chunks.append({
                    "doc_id": doc_id,
                    "file_name": file_name,
                    "page_number": page_number,
                    "chunk_id": chunk_id,
                    "char_start": current_char_start,
                    "char_end": chunk_end,
                    "text": chunk_text,
                })
                
                overlap_words = int(len(current_chunk) * (self.overlap / self.chunk_size))
                if overlap_words > 0:
                    current_chunk = current_chunk[-overlap_words:]
                    overlap_text = ' '.join(current_chunk)
                    current_char_start = chunk_end - len(overlap_text)
                else:
                    current_chunk = []
                    current_char_start = chunk_end + 1
            
            char_position += len(word) + 1
        
        if current_chunk:
            chunk_id = str(uuid.uuid4())
            chunk_text = ' '.join(current_chunk)
            chunks.append({
                "doc_id": doc_id,
                "file_name": file_name,
                "page_number": page_number,
                "chunk_id": chunk_id,
                "char_start": current_char_start,
                "char_end": char_position,
                "text": chunk_text,
            })
        
        return chunks


class FileStorage:
    
    def __init__(self):
        self.courses: Dict[str, Dict] = {}
        self.chunker = TextChunker()
        os.makedirs(PERSIST_UPLOADS, exist_ok=True)
        os.makedirs(PERSIST_COURSES, exist_ok=True)

    def _course_dir(self, course_id: str) -> str:
        return os.path.join(PERSIST_UPLOADS, course_id)

    def _manifest_path(self, course_id: str) -> str:
        return os.path.join(PERSIST_COURSES, f"{course_id}.json")

    def _save_manifest(self, course_id: str, processed_files: List[str], all_chunks: List[Dict]):
        manifest = {
            "course_id": course_id,
            "files": processed_files,
            "chunks": all_chunks,
        }
        with open(self._manifest_path(course_id), "w", encoding="utf-8") as f:
            json.dump(manifest, f)
    
    async def save_and_process_files(
        self, 
        files: List[Tuple[str, bytes, str]]
    ) -> Tuple[str, List[str], List[Dict]]:
        course_id = str(uuid.uuid4())
        processed_files = []
        all_chunks = []
        saved_files: List[Dict] = []
        os.makedirs(self._course_dir(course_id), exist_ok=True)
        
        for file_name, file_content, content_type in files:
            if len(file_content) > MAX_FILE_SIZE:
                continue
            
            doc_id = str(uuid.uuid4())
            
            saved_name = f"{doc_id}_{file_name}"
            temp_path = os.path.join(self._course_dir(course_id), saved_name)
            with open(temp_path, "wb") as f:
                f.write(file_content)
            
            try:
                pages = TextExtractor.extract(temp_path, file_name)
                
                for page in pages:
                    chunks = self.chunker.chunk_text(
                        text=page["text"],
                        doc_id=doc_id,
                        file_name=file_name,
                        page_number=page["page_number"],
                        base_char_offset=page["char_start"],
                    )
                    all_chunks.extend(chunks)
                
                processed_files.append(file_name)
                saved_files.append({"file_name": file_name, "saved_name": saved_name})
            finally:
                # Keep the uploaded files persisted for future reuse
                pass
        
        self.courses[course_id] = {
            "course_id": course_id,
            "files": processed_files,
            "chunks": all_chunks,
            "saved_files": saved_files,
        }
        # include saved files in manifest
        manifest = {
            "course_id": course_id,
            "files": processed_files,
            "chunks": all_chunks,
            "saved_files": saved_files,
        }
        with open(self._manifest_path(course_id), "w", encoding="utf-8") as f:
            json.dump(manifest, f)
        
        return course_id, processed_files, all_chunks
    
    def get_course_chunks(self, course_id: str) -> List[Dict]:
        if course_id in self.courses:
            return self.courses[course_id].get("chunks", [])
        return []
    
    def get_course_info(self, course_id: str) -> Optional[Dict]:
        return self.courses.get(course_id)

    def load_existing_courses(self):
        if not os.path.isdir(PERSIST_COURSES):
            return
        for name in os.listdir(PERSIST_COURSES):
            if not name.endswith(".json"):
                continue
            path = os.path.join(PERSIST_COURSES, name)
            try:
                with open(path, "r", encoding="utf-8") as f:
                    manifest = json.load(f)
                course_id = manifest.get("course_id") or os.path.splitext(name)[0]
                files = manifest.get("files", [])
                chunks = manifest.get("chunks", [])
                saved_files = manifest.get("saved_files", [])
                self.courses[course_id] = {
                    "course_id": course_id,
                    "files": files,
                    "chunks": chunks,
                    "saved_files": saved_files,
                }
            except Exception:
                continue

    def get_saved_files(self, course_id: str) -> List[Dict]:
        info = self.courses.get(course_id)
        if not info:
            return []
        return info.get("saved_files", [])

    def get_file_details(self, course_id: str, saved_name: str) -> Dict:
        info = self.courses.get(course_id) or {}
        saved_files = info.get("saved_files", [])
        original_name = None
        for sf in saved_files:
            if sf.get("saved_name") == saved_name:
                original_name = sf.get("file_name")
                break
        base_dir = self._course_dir(course_id)
        path = os.path.join(base_dir, saved_name)
        size = os.path.getsize(path) if os.path.exists(path) else 0
        chunks = info.get("chunks", [])
        file_chunks = [ch for ch in chunks if ch.get("file_name") == original_name]
        page_count = 0
        if file_chunks:
            try:
                page_count = max(ch.get("page_number", 0) for ch in file_chunks)
            except Exception:
                page_count = 0
        return {
            "file_name": original_name or saved_name,
            "size_bytes": size,
            "page_count": page_count,
            "chunk_count": len(file_chunks),
            "url": f"/files/{course_id}/{saved_name}",
        }

    def delete_file(self, course_id: str, saved_name: str) -> bool:
        info = self.courses.get(course_id)
        if not info:
            return False
        saved_files = info.get("saved_files", [])
        original_name = None
        new_saved = []
        for sf in saved_files:
            if sf.get("saved_name") == saved_name:
                original_name = sf.get("file_name")
            else:
                new_saved.append(sf)
        base_dir = self._course_dir(course_id)
        path = os.path.join(base_dir, saved_name)
        if os.path.exists(path):
            try:
                os.remove(path)
            except Exception:
                pass
        chunks = info.get("chunks", [])
        remaining_chunks = [ch for ch in chunks if ch.get("file_name") != original_name]
        files_list = [fn for fn in info.get("files", []) if fn != original_name]
        self.courses[course_id] = {
            "course_id": course_id,
            "files": files_list,
            "chunks": remaining_chunks,
            "saved_files": new_saved,
        }
        manifest = {
            "course_id": course_id,
            "files": files_list,
            "chunks": remaining_chunks,
            "saved_files": new_saved,
        }
        try:
            with open(self._manifest_path(course_id), "w", encoding="utf-8") as f:
                json.dump(manifest, f)
        except Exception:
            pass
        return True

    def delete_all(self) -> int:
        count = 0
        for course_id, info in list(self.courses.items()):
            base_dir = self._course_dir(course_id)
            try:
                if os.path.isdir(base_dir):
                    for name in os.listdir(base_dir):
                        try:
                            os.remove(os.path.join(base_dir, name))
                        except Exception:
                            pass
                    try:
                        os.rmdir(base_dir)
                    except Exception:
                        pass
            except Exception:
                pass
            manifest_path = self._manifest_path(course_id)
            try:
                if os.path.exists(manifest_path):
                    os.remove(manifest_path)
            except Exception:
                pass
            count += 1
        self.courses = {}
        return count


file_storage = FileStorage()
